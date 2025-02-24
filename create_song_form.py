from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class ChangeLyrics:
    def __init__(self, text, save_path, ppt_title, ppt_path=None):
        self.text = "\n".join(text.splitlines())
        self.ppt_path = ppt_path
        self.save_path = save_path
        self.ppt_title = ppt_title

    def sep_lyrics(self):
        full_lyrics = [line.strip() for line in self.text.strip().split('\n') if line.strip()]
        korean_lines = []
        english_lines = []
        # 처음에는 일반 순서: 첫 줄은 한글, 두 번째 줄은 영어

        i = 0
        while i < len(full_lyrics):
            if full_lyrics[i] == "!empty_slide":
                # marker를 양쪽 리스트에 추가하고, 이후부터 순서를 토글
                korean_lines.append("!empty_slide")
                english_lines.append("!empty_slide")
                i += 1
            else:
                if i + 1 < len(full_lyrics):
                    korean_lines.append(full_lyrics[i])
                    english_lines.append(full_lyrics[i + 1])
                    i += 2
                else:
                    # 남은 라인이 하나라면 상황에 맞게 처리
                    korean_lines.append(full_lyrics[i])
                    english_lines.append(full_lyrics[i])
                    i += 1
        return korean_lines, english_lines

    # 두 줄씩 그룹화
    def group_pairs(self, lines):
        processed_lines = []

        i = 0
        while i < len(lines):
            if lines[i].lower() == "!blank":
                i += 1
                continue

            if lines[i].lower() == "!empty_slide":
                processed_lines.append(["",""])  # 공백 슬라이드 (단독 리스트)
                i += 1
                continue

            if i + 1 < len(lines) and lines[i + 1].lower() not in ["!blank", "!empty_slide"]:
                processed_lines.append([lines[i], lines[i + 1]])  # 두 줄씩 묶음
                i += 2
            else:
                processed_lines.append([lines[i]])  # 한 줄만 있는 경우 단독 리스트
                i += 1
        return processed_lines

    def setting_slide_size_16_9(self, prs):
        prs.slide_width = Inches(13.3333)  # 약 12288000 EMU
        prs.slide_height = Inches(7.5)

    def setting_slide_layout(self, prs):
        # 슬라이드 레이아웃 선택(빈 슬라이드)
        slide_layout = prs.slide_layouts[6]
        background = slide_layout.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0,0,0)
        return slide_layout

    def create_text_box(self, left, top, width, height):
        left = Inches(float(left))
        top = Inches(float(top))
        width = Inches(float(width))
        height = Inches(float(height))

        # txBox = slide.shapes.add_textbox(left, top, width, height)
        return left, top, width, height

    def text_exchange(self, pair, slide_index):
        return "\n".join(pair[slide_index])

    def input_text(self, pair):
        print(pair)
        return "\n".join(pair)

    def font_settings(self, paragraph, font, size):
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = font  # 글씨체
            run.font.size = Pt(int(size))  # 글자 크기
            run.font.bold = True  # 굵게 설정
            run.font.color.rgb = RGBColor(255, 255, 255)  # 글자 색

    def save(self, prs):
        save_path = f"{self.save_path}/{self.ppt_title}.pptx"
        prs.save(save_path)
        return print("파일 저장됨.")

    def change_lyrics(self):
        korean_lines, english_lines = self.sep_lyrics()
        korean_pairs = self.group_pairs(korean_lines)
        english_pairs = self.group_pairs(english_lines)

        # PPT 파일 열기
        prs = Presentation(self.ppt_path)
        self.setting_slide_size_16_9(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, slide in enumerate(prs.slides):
            if slide_index >= len(korean_pairs) or slide_index >= len(english_pairs):
                print(f"슬라이드 {slide_index + 1}에 텍스트가 부족합니다.")
                continue

            # 현재 슬라이드의 텍스트 상자 저장
            text_boxes = [shape.text_frame for shape in slide.shapes if shape.has_text_frame]

            # 텍스트 교환
            korean_text = self.text_exchange(korean_pairs, slide_index)
            english_text = self.text_exchange(english_pairs, slide_index)
            text_boxes[0].text = korean_text
            text_boxes[1].text = english_text
            # 텍스트 상자별로 처리
            for text_box_index, text_frame in enumerate(text_boxes):
                # 첫 번째 텍스트 상자 처리
                if text_box_index == 0:
                    for paragraph in text_frame.paragraphs:
                        self.font_settings(paragraph, "나눔고딕", "40")
                else:
                    for paragraph in text_frame.paragraphs:
                        self.font_settings(paragraph, "나눔고딕", "15")

        self.save(prs)

    def change_only_korean_lyrics(self):
        # 텍스트를 줄 단위로 분리하고 빈 줄(공백) 제거
        korean_lines = [line.strip() for line in self.text.strip().split('\n') if line.strip()]
        # 텍스트를 두개씩 하나의 그룹으로 형성
        korean_pairs = self.group_pairs(korean_lines)

        # PPT 생성
        prs = Presentation(self.ppt_path)
        self.setting_slide_size_16_9(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, slide in enumerate(prs.slides):
            if slide_index >= len(korean_pairs):
                print(f"슬라이드 {slide_index + 1}에 텍스트가 부족합니다.")
                continue

            # 현재 슬라이드의 텍스트 상자 저장
            text_boxes = [shape.text_frame for shape in slide.shapes if shape.has_text_frame]

            # 텍스트 교환
            korean_text = self.text_exchange(korean_pairs, slide_index)
            text_boxes[0].text = korean_text

            # 텍스트 설정
            for text_box_index, text_frame in enumerate(text_boxes):
                if text_box_index == 0:
                    for paragraph in text_frame.paragraphs:
                        self.font_settings(paragraph, "나눔고딕", "40")

        self.save(prs)

    def change_only_english_lyrics(self):
        # 텍스트를 줄 단위로 분리하고 빈 줄(공백) 제거
        english_lines = [line.strip() for line in self.text.strip().split('\n') if line.strip()]
        # 텍스트를 두개씩 하나의 그룹으로 형성
        english_pairs = self.group_pairs(english_lines)

        # PPT 생성
        prs = Presentation(self.ppt_path)
        self.setting_slide_size_16_9(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, slide in enumerate(prs.slides):
            if slide_index >= len(english_pairs):
                print(f"슬라이드 {slide_index + 1}에 텍스트가 부족합니다.")
                continue

            # 현재 슬라이드의 텍스트 상자 저장
            text_boxes = [shape.text_frame for shape in slide.shapes if shape.has_text_frame]

            # 텍스트 교환
            english_text = self.text_exchange(english_pairs, slide_index)
            text_boxes[1].text = english_text

            # 텍스트 설정
            for text_box_index, text_frame in enumerate(text_boxes):
                if text_box_index == 1:
                    for paragraph in text_frame.paragraphs:
                        self.font_settings(paragraph, "나눔고딕", "15")

        self.save(prs)

    def create_lyrics(self):
        # 한글과 영어 라인 분리 (홀수 인덱스: 한글, 짝수 인덱스: 영어)
        korean_lines, english_lines = self.sep_lyrics()
        korean_pairs = self.group_pairs(korean_lines)
        english_pairs = self.group_pairs(english_lines)
        # PPT 생성
        prs = Presentation()
        self.setting_slide_size_16_9(prs)
        slide_layout = self.setting_slide_layout(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, (ko_pair, en_pair) in enumerate(zip(korean_pairs, english_pairs), start=1):

            # 빈 슬라이드 생성
            slide = prs.slides.add_slide(slide_layout)

            # 슬라이드 배경을 검은색으로 설정
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # 텍스트 박스 1: 한글 텍스트
            left_ko, top_ko, width_ko, height_ko = self.create_text_box(1.6667, 2.55, 10, 1.476)
            # 텍스트 박스 2: 영어 텍스트
            left_en, top_en, width_en, height_en = self.create_text_box(1.6667, 4.572, 10, 0.933)

            # # 슬라이드에 텍스트 박스 추가
            txBox_ko = slide.shapes.add_textbox(left_ko, top_ko, width_ko, height_ko)
            txBox_en = slide.shapes.add_textbox(left_en, top_en, width_en, height_en)

            # 텍스트 입력
            korean_text = self.input_text(ko_pair)
            english_text = self.input_text(en_pair)
            txBox_ko.text = korean_text
            txBox_en.text = english_text

            # 텍스트 설정
            for paragraph in txBox_ko.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕", "40")

            for paragraph in txBox_en.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕", "15")

        self.save(prs)

    def create_only_korean_lyrics(self):
        # 텍스트를 줄 단위로 분리하고 빈 줄(공백) 제거
        korean_lines = [line.strip() for line in self.text.strip().split('\n') if line.strip()]
        # 텍스트를 두개씩 하나의 그룹으로 형성
        korean_pairs = self.group_pairs(korean_lines)

        # PPT 생성
        prs = Presentation()
        self.setting_slide_size_16_9(prs)
        slide_layout = self.setting_slide_layout(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, ko_pair in enumerate(korean_pairs, start=1):

            # 빈 슬라이드 생성
            slide = prs.slides.add_slide(slide_layout)

            # 슬라이드 배경을 검은색으로 설정
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # 텍스트 박스 1: 한글 텍스트
            left_ko, top_ko, width_ko, height_ko = self.create_text_box(1.6667, 2.55, 10, 1.476)
            # 텍스트 박스 2: 영어 텍스트
            left_en, top_en, width_en, height_en = self.create_text_box(1.6667, 4.572, 10, 0.933)

            # # 슬라이드에 텍스트 박스 추가
            txBox_ko = slide.shapes.add_textbox(left_ko, top_ko, width_ko, height_ko)
            txBox_en = slide.shapes.add_textbox(left_en, top_en, width_en, height_en)

            # 텍스트 입력
            korean_text = self.input_text(ko_pair)
            txBox_ko.text = korean_text

            # 텍스트 설정
            for paragraph in txBox_ko.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕", "40")

            for paragraph in txBox_en.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕", "15")

        self.save(prs)

    def create_only_english_lyrics(self):
        # 텍스트를 줄 단위로 분리하고 빈 줄(공백) 제거
        english_lines = [line.strip() for line in self.text.strip().split('\n') if line.strip()]
        # 텍스트를 두개씩 하나의 그룹으로 형성
        english_pairs = self.group_pairs(english_lines)

        # PPT 생성
        prs = Presentation()
        self.setting_slide_size_16_9(prs)
        slide_layout = self.setting_slide_layout(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, en_pair in enumerate(english_pairs, start=1):

            # 빈 슬라이드 생성
            slide = prs.slides.add_slide(slide_layout)

            # 슬라이드 배경을 검은색으로 설정
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # 텍스트 박스 1: 한글 텍스트
            left_ko, top_ko, width_ko, height_ko = self.create_text_box(1.6667, 2.55, 10, 1.476)
            # 텍스트 박스 2: 영어 텍스트
            left_en, top_en, width_en, height_en = self.create_text_box(1.6667, 4.572, 10, 0.933)

            # # 슬라이드에 텍스트 박스 추가
            txBox_ko = slide.shapes.add_textbox(left_ko, top_ko, width_ko, height_ko)
            txBox_en = slide.shapes.add_textbox(left_en, top_en, width_en, height_en)

            # 텍스트 입력
            english_text = self.input_text(en_pair)
            txBox_en.text = english_text

            # 텍스트 설정
            for paragraph in txBox_ko.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕", "40")

            for paragraph in txBox_en.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕", "15")

        self.save(prs)

    def create_seoul_form(self):
        # 한글과 영어 라인 분리 (홀수 인덱스: 한글, 짝수 인덱스: 영어)
        korean_lines, english_lines = self.sep_lyrics()
        korean_pairs = self.group_pairs(korean_lines)
        english_pairs = self.group_pairs(english_lines)
        # PPT 생성
        prs = Presentation()
        self.setting_slide_size_16_9(prs)
        slide_layout = self.setting_slide_layout(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, (ko_pair, en_pair) in enumerate(zip(korean_pairs, english_pairs), start=1):

            # 빈 슬라이드 생성
            slide = prs.slides.add_slide(slide_layout)

            # 슬라이드 배경을 검은색으로 설정
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # 텍스트 박스 1: 한글 텍스트
            left_ko, top_ko, width_ko, height_ko = self.create_text_box(0.8417, 0.5723, 11.6538, 1.4818)
            # 텍스트 박스 2: 영어 텍스트
            left_en, top_en, width_en, height_en = self.create_text_box(0.3111, 2.0536, 12.7131, 1.0432)

            # # 슬라이드에 텍스트 박스 추가
            txBox_ko = slide.shapes.add_textbox(left_ko, top_ko, width_ko, height_ko)
            txBox_en = slide.shapes.add_textbox(left_en, top_en, width_en, height_en)

            # 텍스트 입력
            korean_text = self.input_text(ko_pair)
            english_text = self.input_text(en_pair)
            txBox_ko.text = korean_text
            txBox_en.text = english_text

            # 텍스트 설정
            for paragraph in txBox_ko.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕OTF ExtraBold", "41")

            for paragraph in txBox_en.text_frame.paragraphs:
                self.font_settings(paragraph, "나눔고딕OTF ExtraBold", "28")

        self.save(prs)

    def create_subtitle_lyrics(self):
        # 한글과 영어 라인 분리 (홀수 인덱스: 한글, 짝수 인덱스: 영어)
        korean_lines, english_lines = self.sep_lyrics()
        # PPT 생성
        prs = Presentation()
        self.setting_slide_size_16_9(prs)
        slide_layout = self.setting_slide_layout(prs)

        # 슬라이드별로 텍스트 상자를 순회하며 처리
        for slide_index, (ko_line, en_line) in enumerate(zip(korean_lines, english_lines), start=1):

            # 빈 슬라이드 생성
            slide = prs.slides.add_slide(slide_layout)

            # 슬라이드 배경을 검은색으로 설정
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # 텍스트 박스 1: 한글 텍스트
            left_ko, top_ko, width_ko, height_ko = self.create_text_box(1.6667, 6.17, 10, 1.476)
            # 텍스트 박스 2: 영어 텍스트
            left_en, top_en, width_en, height_en = self.create_text_box(1.6667, 7.03, 10, 0.933)

            # # 슬라이드에 텍스트 박스 추가
            txBox_ko = slide.shapes.add_textbox(left_ko, top_ko, width_ko, height_ko)
            txBox_en = slide.shapes.add_textbox(left_en, top_en, width_en, height_en)

            # 텍스트 입력
            if "!empty_slide" in ko_line or "!empty_slide" in en_line:
                continue
            else:
                txBox_ko.text = ko_line
                txBox_en.text = en_line

                # 텍스트 설정
                for paragraph in txBox_ko.text_frame.paragraphs:
                    self.font_settings(paragraph, "나눔고딕", "40")

                for paragraph in txBox_en.text_frame.paragraphs:
                    self.font_settings(paragraph, "나눔고딕", "15")

        self.save(prs)

    def assemble_song(verse_names, verse_texts, order):
        """
        verse_names: 리스트, 예를 들어 ["verse1", "verse2", "chorus"]
        verse_texts: 리스트, 각 구절에 해당하는 텍스트
        order: 문자열, 예: "verse1 - verse2 - chorus - verse3 - chorus - chorus"
        """
        # 구절 이름과 가사를 딕셔너리로 매핑
        verses = {}
        for name, text in zip(verse_names, verse_texts):
            verses[name.strip()] = text.strip()

        # order 문자열을 '-' 기준으로 분리하여 순서대로 배열 생성
        order_list = [item.strip() for item in order.split('-') if item.strip()]

        # 지정한 순서대로 구절을 조합
        assembled_song = []
        for item in order_list:
            if item in verses:
                assembled_song.append(verses[item])
            else:
                assembled_song.append(f"존재하지 않는 구절: {item}")

        # 최종 텍스트 반환
        final_text = "\n\n".join(assembled_song)
        return final_text