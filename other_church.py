import os
import logging
import re
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# 로깅 설정
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class ChangeBibleFormat2:
    def __init__(self, text, bible_path, output_path, ppt_title):
        self.text = "\n".join(text.splitlines())
        self.book_names = {
            "창": "창세기", "출": "출애굽기", "레": "레위기", "민": "민수기", "신": "신명기",
            "수": "여호수아", "삿": "사사기", "룻": "룻기", "삼상": "사무엘상", "삼하": "사무엘하",
            "왕상": "열왕기상", "왕하": "열왕기하", "대상": "역대상", "대하": "역대하", "스": "에스라",
            "느": "느헤미야", "에": "에스더", "욥": "욥기", "시": "시편", "잠": "잠언", "전": "전도서",
            "아": "아가", "사": "이사야", "렘": "예레미야", "애": "예레미야애가", "겔": "에스겔",
            "단": "다니엘", "호": "호세아", "욜": "요엘", "암": "아모스", "옵": "오바댜", "욘": "요나",
            "미": "미가", "나": "나훔", "합": "하박국", "습": "스바냐", "학": "학개", "슥": "스가랴",
            "말": "말라기", "마": "마태복음", "막": "마가복음", "눅": "누가복음", "요": "요한복음",
            "행": "사도행전", "롬": "로마서", "고전": "고린도전서", "고후": "고린도후서", "갈": "갈라디아서",
            "엡": "에베소서", "빌": "빌립보서", "골": "골로새서", "살전": "데살로니가전서",
            "살후": "데살로니가후서", "딤전": "디모데전서", "딤후": "디모데후서", "딛": "디도서",
            "몬": "빌레몬서", "히": "히브리서", "약": "야고보서", "벧전": "베드로전서",
            "벧후": "베드로후서", "요일": "요한일서", "요이": "요한이서", "요삼": "요한삼서",
            "유": "유다서", "계": "요한계시록"
        }
        self.patterns = {
            "main_verses": r"(.+?)\s(\d+)장\s(\d+)~(\d+)절",
            "single_verse": r"(.+?)\s(\d+):(\d+)",
            "single_verse_with_comma": r"(.+?)\s(\d+):(\d+),(\d+)",
            "single_verse_with_verse_range": r"(.+?)\s(\d+):(\d+),(\d+)~(\d+)",
            "verse_range": r"(.+?)\s(\d+):(\d+)~(\d+)",
            "verse_range_with_comma": r"(.+?)\s(\d+):(\d+)~(\d+),(\d+)",
            "complex_range_with_comma": r"(.+?)\s(\d+):(\d+)~(\d+),(\d+)~(\d+)"
        }
        self.bible_names = {
            "창세기":"구약", "출애굽기":"구약", "레위기":"구약",
            "민수기":"구약", "신명기":"구약", "여호수아":"구약",
            "사사기":"구약", "룻기":"구약", "사무엘상":"구약",
            "사무엘하":"구약", "열왕기상":"구약", "열왕기하":"구약",
            "역대상":"구약", "역대하":"구약", "에스라":"구약",
            "느헤미야":"구약", "에스더":"구약", "욥기":"구약",
            "시편":"구약", "잠언":"구약", "전도서":"구약",
            "아가":"구약", "이사야":"구약", "예레미야":"구약",
            "예레미야애가":"구약", "에스겔":"구약", "다니엘":"구약",
            "호세아":"구약", "요엘":"구약", "아모스":"구약",
            "오바댜":"구약", "요나":"구약", "미가":"구약",
            "나훔":"구약", "하박국":"구약", "스바냐":"구약",
            "학개":"구약", "스가랴":"구약", "말라기":"구약",
            "마태복음":"신약", "마가복음":"신약", "누가복음":"신약",
            "요한복음":"신약", "사도행전":"신약", "로마서":"신약",
            "고린도전서":"신약", "고린도후서":"신약", "갈라디아서":"신약",
            "에베소서":"신약", "빌립보서":"신약", "골로새서":"신약",
            "데살로니가전서":"신약", "데살로니가후서":"신약", "디모데전서":"신약",
            "디모데후서":"신약", "디도서":"신약", "빌레몬서":"신약",
            "히브리서":"신약", "야고보서":"신약", "베드로전서":"신약",
            "베드로후서":"신약", "요한일서":"신약", "요한이서":"신약",
            "요한삼서":"신약", "유다서":"신약", "요한계시록":"신약"
        }
        self.bible_path = bible_path
        self.output_path = output_path
        self.ppt_title = ppt_title

    def change_book_names_to_path(self, book):
        for book_name, testament in self.bible_names.items():
            if book == book_name:
                path_without_bible_name = os.path.join(self.bible_path, testament)
                if not os.path.isdir(path_without_bible_name):
                    logging.error(f"테스먼트 폴더가 존재하지 않습니다: {path_without_bible_name}")
                    return None
                for subfolder in os.listdir(path_without_bible_name):
                    if subfolder.endswith(book):
                        full_path = os.path.join(testament, subfolder)
                        logging.info(f"Book path found: {full_path}")
                        return full_path
        logging.error(f"책 이름을 경로로 변경할 수 없습니다: {book}")
        return None

    def create_ppt_path_and_full_name(self, book, chapter, expanded_verse):
        book_path = self.change_book_names_to_path(book)
        if not book_path:
            raise ValueError(f"책 경로를 찾을 수 없습니다: {book}")
        ppt_file_path = os.path.join(self.bible_path, book_path, f"{book}{chapter}장.pptx")
        split_name = expanded_verse.split(" ")
        full_bible_name = ""
        for korean_name, english_name in self.english_book_names.items():
            if split_name[0] == korean_name:
                full_bible_name = f"{split_name[0]} {english_name} | {' '.join(split_name[1:])}"
                break
        return ppt_file_path, full_bible_name

    def spacing_name(self):
        find_pattern = r'[가-힣]+\s?\d+장?\s?\d+(?:~|-)\d+절?|[가-힣]+\s?\d+:\d+(?:~\d+(?:,\d+(?:~\d+)?)?|,\d+(?:~\d+)?)?'
        sub_pattern = r'([가-힣]+)\s?(\d+:\d+(?:~\d+(?:,\d+(?:~\d+)?)?|,\d+(?:~\d+)?)?)'
        replacement = r'\1 \2'
        spacing_verses = []
        verses = re.findall(find_pattern, self.text)
        for verse in verses:
            spacing = re.sub(sub_pattern, replacement, verse)
            spacing_verses.append(spacing)
        return spacing_verses

    def expanding_name(self):
        spacing_verses = self.spacing_name()
        expanded_verses = []
        for verse in spacing_verses:
            replaced = False
            book = verse.split(" ")[0]
            for abbrev, full_name in self.book_names.items():
                if book == abbrev:
                    expanded = verse.replace(abbrev, full_name, 1)
                    if len(expanded) > len(full_name) and expanded[len(full_name)] != ' ':
                        expanded = full_name + ' ' + expanded[len(full_name):]
                    expanded_verses.append(expanded)
                    replaced = True
                    break
                elif book == full_name:
                    expanded = verse
                    if len(expanded) > len(full_name) and expanded[len(full_name)] != ' ':
                        expanded = full_name + ' ' + expanded[len(full_name):]
                    expanded_verses.append(expanded)
                    replaced = True
                    break
            if not replaced:
                raise ValueError(f"Unknown name : {verse}")
        return expanded_verses

    def create_ppt_path_and_slide(self):
        path_list = []
        s_range = []
        expanded_verses = self.expanding_name()
        for idx, expanded_verse in enumerate(expanded_verses):
            if re.match(self.patterns.get("complex_range_with_comma"), expanded_verse):
                match = re.match(self.patterns.get("complex_range_with_comma"), expanded_verse)
                book, chapter, start_verse, end_verse, addition_start_verse, addition_end_verse = match.groups()
                start_verse, end_verse, addition_start_verse, addition_end_verse = int(start_verse), int(end_verse), int(addition_start_verse), int(addition_end_verse)
                slide_range = list(range(start_verse, end_verse + 1)) + list(range(addition_start_verse, addition_end_verse + 1))
            elif re.match(self.patterns.get("verse_range_with_comma"), expanded_verse):
                match = re.match(self.patterns.get("verse_range_with_comma"), expanded_verse)
                book, chapter, start_verse, end_verse, addition_verse = match.groups()
                start_verse, end_verse, addition_verse = int(start_verse), int(end_verse), int(addition_verse)
                slide_range = list(range(start_verse, end_verse + 1)) + [addition_verse]
            elif re.match(self.patterns.get("verse_range"), expanded_verse):
                match = re.match(self.patterns.get("verse_range"), expanded_verse)
                book, chapter, start_verse, end_verse = match.groups()
                start_verse, end_verse = int(start_verse), int(end_verse)
                slide_range = list(range(start_verse, end_verse + 1))
            elif re.match(self.patterns.get("single_verse_with_verse_range"), expanded_verse):
                match = re.match(self.patterns.get("single_verse_with_verse_range"), expanded_verse)
                book, chapter, verse, addition_start_verse, addition_end_verse = match.groups()
                verse, addition_start_verse, addition_end_verse = int(verse), int(addition_start_verse), int(addition_end_verse)
                slide_range = [verse] + list(range(addition_start_verse, addition_end_verse + 1))
            elif re.match(self.patterns.get("single_verse_with_comma"), expanded_verse):
                match = re.match(self.patterns.get("single_verse_with_comma"), expanded_verse)
                book, chapter, verse, addition_verse = match.groups()
                verse, addition_verse = int(verse), int(addition_verse)
                slide_range = [verse, addition_verse]
            elif re.match(self.patterns.get("single_verse"), expanded_verse):
                match = re.match(self.patterns.get("single_verse"), expanded_verse)
                book, chapter, start_verse = match.groups()
                start_verse = int(start_verse)
                slide_range = [start_verse]
            elif re.match(self.patterns.get("main_verses"), expanded_verse):
                match = re.match(self.patterns.get("main_verses"), expanded_verse)
                book, chapter, start_verse, end_verse = match.groups()
                start_verse, end_verse = int(start_verse), int(end_verse)
                slide_range = list(range(start_verse, end_verse + 1))
                expanded_verses[idx] = f"{book} {chapter}:{start_verse}~{end_verse}"
            else:
                raise ValueError("Invalid verse format. Use 'Book Chapter:Verse-Verse'.")

            book_path = self.change_book_names_to_path(book)
            ppt_file_path = os.path.join(self.bible_path, book_path, f"{book}{chapter}장.pptx")
            if not os.path.exists(ppt_file_path):
                logging.error(f"PPT 파일이 존재하지 않습니다: {ppt_file_path}")
                raise FileNotFoundError(f"PPT 파일이 존재하지 않습니다: {ppt_file_path}")
            path_list.append(ppt_file_path)
            s_range.append(slide_range)

        return path_list, s_range, expanded_verses

    def copy_slide(self, prs, slide, full_bible_name):
        """슬라이드의 모든 요소를 새로운 슬라이드로 복사하고 서식을 적용."""
        # 빈 슬라이드 추가
        slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        new_slide = prs.slides.add_slide(slide_layout)

        # 슬라이드의 모든 도형을 복사
        for idx, shape in enumerate(slide.shapes):
            if idx == 0:
                left_title = Inches(4.76)
                top_title = Inches(0.46)
                width_title = Inches(3.81)
                height_title = Inches(0.91)
                txBox_title = new_slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
                text_frame = txBox_title.text_frame

                # 첫 번째 텍스트 상자에만 full_bible_name 설정
                text_frame.text = full_bible_name
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.underline = True # 밑줄
                p.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.name = '맑은 고딕'
            elif idx == 1:
                left_content = Inches(0.93)
                top_content = Inches(1.78)
                width_content = Inches(11.48)
                height_content = Inches(4.14)
                txBox_content = new_slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
                text_frame = txBox_content.text_frame
                # 다른 텍스트 상자 복사
                text_frame.text = shape.text_frame.text
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.name = '맑은 고딕'

    def create_ppt_file(self):
        try:
            path_list, s_range, expanded_verses = self.create_ppt_path_and_slide()
            merged_ppt = Presentation()
            # 슬라이드 크기 설정 (필요에 따라 조정)
            merged_ppt.slide_width = Inches(13.3333)
            merged_ppt.slide_height = Inches(7.5)

            for ppt_path, slide_range, expanded_verse in zip(path_list, s_range, expanded_verses):
                source_prs = Presentation(ppt_path)

                for slide_num in slide_range:
                    if slide_num < 1 or slide_num > len(source_prs.slides):
                        logging.warning(f"슬라이드 번호 {slide_num}이(가) {ppt_path}에 없습니다.")
                        continue
                    source_slide = source_prs.slides[slide_num - 1]
                    try:
                        self.copy_slide(merged_ppt, source_slide, expanded_verse)
                        logging.info(f"슬라이드 {slide_num}을(를) 성공적으로 복사했습니다.")
                    except Exception as e:
                        logging.error(f"슬라이드 {slide_num} 복사 중 오류 발생: {e}")

            # 모든 슬라이드의 배경 색상 변경
            for slide in merged_ppt.slides:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색

            # 최종 저장
            save_path = os.path.join(self.output_path, f"{self.ppt_title}.pptx")
            merged_ppt.save(save_path)
            logging.info(f"PPT가 성공적으로 저장되었습니다: {save_path}")

        except Exception as e:
            logging.error(f"에러 발생: {e}")
            raise e
